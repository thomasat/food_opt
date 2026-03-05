"""Generate a PowerPoint slide deck summarizing the Formulation Assistant."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Theme colors ---
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x50)
MID_BLUE = RGBColor(0x34, 0x49, 0x5E)
ACCENT = RGBColor(0x27, 0xAE, 0x60)  # Green accent
LIGHT_BG = RGBColor(0xEC, 0xF0, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
BODY_TEXT = RGBColor(0x4A, 0x4A, 0x4A)
TABLE_HEADER_BG = RGBColor(0x2C, 0x3E, 0x50)
TABLE_ALT_BG = RGBColor(0xF7, 0xF9, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, top=Inches(0), height=Inches(0.08)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_footer(slide, text="Formulation Assistant"):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def add_title_text(slide, text, left, top, width, height, size=36, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_body_text(slide, text, left, top, width, height, size=18, color=BODY_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return tf


def add_bullet_list(slide, items, left, top, width, height, size=18, color=BODY_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.level = 0
    return tf


def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.45)):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    col_width = width // n_cols
    table_height = row_height * n_rows
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height)
    table = shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = BODY_TEXT

    return table


# =============================================
# SLIDE 1: Title
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# Accent bar at top
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_title_text(slide, "Formulation Assistant",
               Inches(1), Inches(2), Inches(11), Inches(1.2),
               size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_title_text(slide, "AI-Powered Food Product Optimization",
               Inches(1), Inches(3.3), Inches(11), Inches(0.8),
               size=28, color=RGBColor(0xBD, 0xC3, 0xC7), bold=False, alignment=PP_ALIGN.CENTER)
add_title_text(slide, "Bayesian Optimization for Ingredient Formulation",
               Inches(1), Inches(4.2), Inches(11), Inches(0.6),
               size=20, color=RGBColor(0x95, 0xA5, 0xA6), bold=False, alignment=PP_ALIGN.CENTER)

# =============================================
# SLIDE 2: The Challenge
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "The Challenge", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

items = [
    "Developing new food formulations requires many trial-and-error iterations",
    "Each lab trial is expensive and time-consuming",
    "Sensory panel evaluations add further cost per iteration",
    "Traditional approaches (DOE, one-variable-at-a-time) scale poorly with many ingredients",
    "Interactions between ingredients are hard to predict manually",
]
add_bullet_list(slide, items, Inches(1.0), Inches(1.8), Inches(7), Inches(4), size=20)

# Callout box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.2), Inches(11), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
shape.line.color.rgb = ACCENT
txf = shape.text_frame
txf.word_wrap = True
txf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = txf.paragraphs[0]
p.text = "Goal: Minimize the number of lab experiments needed to find the optimal recipe."
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = DARK_BLUE

# =============================================
# SLIDE 3: Our Solution
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "Our Solution: Bayesian Optimization", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

add_body_text(slide, "A smart experimental design system that:", Inches(1.0), Inches(1.6), Inches(10), Inches(0.5), size=22, color=DARK_BLUE)

items = [
    "Learns from every experiment you run",
    "Models the relationship between ingredients and sensory outcomes",
    "Suggests the most informative next experiment to try",
    "Converges on the optimal formulation in far fewer iterations",
]
add_bullet_list(slide, items, Inches(1.2), Inches(2.3), Inches(6), Inches(3), size=20)

# Comparison callout
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(2.0), Inches(4.5), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BG
shape.line.color.rgb = MID_BLUE
txf = shape.text_frame
txf.word_wrap = True
p = txf.paragraphs[0]
p.text = "Typical Convergence"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = DARK_BLUE
p.alignment = PP_ALIGN.CENTER
p = txf.add_paragraph()
p.text = ""
p = txf.add_paragraph()
p.text = "Traditional DOE:"
p.font.size = Pt(16)
p.font.color.rgb = BODY_TEXT
p = txf.add_paragraph()
p.text = "50-100+ experiments"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0xE7, 0x4C, 0x3C)
p.font.bold = True
p = txf.add_paragraph()
p.text = ""
p = txf.add_paragraph()
p.text = "Bayesian Optimization:"
p.font.size = Pt(16)
p.font.color.rgb = BODY_TEXT
p = txf.add_paragraph()
p.text = "10-15 experiments"
p.font.size = Pt(16)
p.font.color.rgb = ACCENT
p.font.bold = True

# =============================================
# SLIDE 4: How It Works - The Ask-Tell Loop
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "How It Works: The Ask-Tell Loop", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

steps = [
    ("1", "Setup", "Define ingredients,\nranges & objectives"),
    ("2", "Ask", "System suggests a\nbatch of recipes"),
    ("3", "Lab Trial", "Prepare & evaluate\nthe recipes"),
    ("4", "Tell", "Input sensory panel\nratings into system"),
    ("5", "Repeat", "System learns &\nsuggests improvements"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.8 + i * 2.5)
    y = Inches(2.2)

    # Circle with number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), y, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT if i < 4 else DARK_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_title_text(slide, title, x, y + Inches(1.1), Inches(2), Inches(0.5),
                   size=20, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_body_text(slide, desc, x, y + Inches(1.6), Inches(2), Inches(1),
                  size=14, color=BODY_TEXT)

    # Arrow between steps
    if i < 4:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(2.1), y + Inches(0.2), Inches(0.4), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
        arrow.line.fill.background()

# Loop-back note
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(5.2), Inches(7), Inches(0.8))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
shape.line.color.rgb = RGBColor(0xE6, 0x7E, 0x22)
txf = shape.text_frame
txf.word_wrap = True
txf.paragraphs[0].text = "Each iteration makes the model smarter. The system balances exploring new regions vs. refining promising ones."
txf.paragraphs[0].font.size = Pt(16)
txf.paragraphs[0].font.color.rgb = BODY_TEXT
txf.paragraphs[0].alignment = PP_ALIGN.CENTER

# =============================================
# SLIDE 5: Setup - Ingredients
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "Setup: Ingredients", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

add_body_text(slide, "Upload your ingredient list via CSV with allowed ranges:",
              Inches(1.0), Inches(1.5), Inches(10), Inches(0.5), size=20, color=DARK_BLUE)

add_table(slide,
    ["Name", "Min", "Max", "Cost", "Protein"],
    [
        ["Flour", "30", "60", "0.50", "10.3"],
        ["Sugar", "5", "25", "0.80", "0.0"],
        ["Butter", "10", "30", "2.10", "0.9"],
        ["Eggs", "5", "15", "1.50", "12.6"],
    ],
    Inches(1.0), Inches(2.3), Inches(7), Inches(0.4)
)

items = [
    "Min/Max: allowed quantity range per ingredient",
    "Additional columns (Cost, Protein, etc.) enable property-based constraints",
    "Process parameters (baking temp, mixing time) can also be optimized jointly",
]
add_bullet_list(slide, items, Inches(1.0), Inches(4.8), Inches(10), Inches(2.5), size=18)

# =============================================
# SLIDE 6: Setup - Objectives
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "Setup: Objectives", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

add_body_text(slide, "Define what \"good\" means for your product:",
              Inches(1.0), Inches(1.5), Inches(10), Inches(0.5), size=20, color=DARK_BLUE)

add_table(slide,
    ["Metric", "Goal", "Weight", "Range"],
    [
        ["Chewiness", "Maximize", "0.6", "0 - 10"],
        ["Sweetness", "Target = 5", "0.3", "0 - 10"],
        ["Graininess", "Minimize", "0.1", "0 - 10"],
    ],
    Inches(1.0), Inches(2.3), Inches(7), Inches(0.45)
)

items = [
    "Maximize, Minimize, or hit a specific Target for each sensory attribute",
    "Weights reflect relative importance of each attribute",
    "Ranges normalize all scores to a common 0-1 scale for fair comparison",
    "Combined into a single Utility Score that the optimizer maximizes",
]
add_bullet_list(slide, items, Inches(1.0), Inches(4.5), Inches(10), Inches(2.5), size=18)

# =============================================
# SLIDE 7: Setup - Constraints
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "Setup: Constraints", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

# Property constraints section
add_title_text(slide, "Property Constraints", Inches(1.0), Inches(1.5), Inches(5), Inches(0.5),
               size=24, color=ACCENT, bold=True)
add_body_text(slide, "Based on ingredient properties from your CSV:",
              Inches(1.0), Inches(2.1), Inches(5), Inches(0.4), size=16, color=BODY_TEXT)
items = [
    "Total Cost <= 5.00 per unit",
    "Total Protein >= 8g per serving",
]
add_bullet_list(slide, items, Inches(1.2), Inches(2.6), Inches(5), Inches(1.5), size=18)

# Quantity constraints section
add_title_text(slide, "Ingredient Quantity Constraints", Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
               size=24, color=ACCENT, bold=True)
add_body_text(slide, "Direct limits on ingredient amounts:",
              Inches(7.0), Inches(2.1), Inches(5.5), Inches(0.4), size=16, color=BODY_TEXT)
items = [
    "Sugar + Honey <= 50g",
    "Total recipe mass: 90g - 110g",
]
add_bullet_list(slide, items, Inches(7.2), Inches(2.6), Inches(5), Inches(1.5), size=18)

# Bottom callout
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.8), Inches(11), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
shape.line.color.rgb = ACCENT
txf = shape.text_frame
txf.word_wrap = True
txf.paragraphs[0].text = "All constraints are enforced during optimization - every suggested recipe is guaranteed to be feasible."
txf.paragraphs[0].font.size = Pt(18)
txf.paragraphs[0].font.bold = True
txf.paragraphs[0].font.color.rgb = DARK_BLUE
txf.paragraphs[0].alignment = PP_ALIGN.CENTER

# =============================================
# SLIDE 8: The Optimization Engine
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "The Optimization Engine", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

# Phase 1
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BG
shape.line.color.rgb = MID_BLUE

add_title_text(slide, "Phase 1: Exploration (First 5 Experiments)",
               Inches(1.1), Inches(1.7), Inches(5), Inches(0.5), size=20, color=DARK_BLUE)
items = [
    "Sobol quasi-random sequences for space-filling design",
    "Ensures broad coverage of the ingredient space",
    "All constraints respected from the start",
    "Optional: pre-screening model to prioritize promising regions",
]
add_bullet_list(slide, items, Inches(1.3), Inches(2.3), Inches(4.8), Inches(1.5), size=15)

# Phase 2
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = LIGHT_BG
shape.line.color.rgb = MID_BLUE

add_title_text(slide, "Phase 2: Optimization (Experiment 6+)",
               Inches(7.3), Inches(1.7), Inches(5), Inches(0.5), size=20, color=DARK_BLUE)
items = [
    "Gaussian Process (GP) models all data collected so far",
    "Noisy Expected Improvement acquisition function",
    "Balances exploration vs. exploitation automatically",
    "Batch suggestions for parallel lab trials",
]
add_bullet_list(slide, items, Inches(7.5), Inches(2.3), Inches(4.8), Inches(1.5), size=15)

# Two modes
add_title_text(slide, "Two Operating Modes", Inches(0.8), Inches(4.6), Inches(11), Inches(0.5),
               size=24, color=DARK_BLUE)

add_table(slide,
    ["Mode", "Best For", "Behavior"],
    [
        ["Standard (Default)", "Smooth synergistic effects", "Faster convergence, assumes gradual interactions"],
        ["Robust", "Sharp thresholds / phase transitions", "Input Warping for non-linear effects, more resilient"],
    ],
    Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.42)
)

# =============================================
# SLIDE 9: Utility Score
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "Utility Score Calculation", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

add_body_text(slide, "Each experiment receives a single score combining all objectives:",
              Inches(1.0), Inches(1.5), Inches(10), Inches(0.5), size=20, color=DARK_BLUE)

steps_text = [
    "1. Normalize each metric to [0, 1] using the defined range",
    "2. Apply goal transformation:  Maximize = score,  Minimize = 1 - score,  Target = 1 - (score - target)^2",
    "3. Weighted sum:  Total Utility = SUM( weight_i  x  utility_i )",
]
add_bullet_list(slide, steps_text, Inches(1.0), Inches(2.2), Inches(10), Inches(2), size=18)

# Example box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(11), Inches(2.3))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
shape.line.color.rgb = RGBColor(0xE6, 0x7E, 0x22)

add_title_text(slide, "Example", Inches(1.3), Inches(4.3), Inches(3), Inches(0.4),
               size=18, color=RGBColor(0xE6, 0x7E, 0x22), bold=True)

example_lines = [
    "Chewiness = 8  (goal=max, weight=0.6, range 0-10)",
    "    normalized = 0.8   |   utility = 0.8   |   weighted = 0.48",
    "",
    "Sweetness = 6  (goal=target 5, weight=0.4, range 0-10)",
    "    normalized = 0.6, target = 0.5   |   utility = 1 - 0.01 = 0.99   |   weighted = 0.396",
    "",
    "Total Utility = 0.48 + 0.396 = 0.876",
]
add_bullet_list(slide, example_lines, Inches(1.3), Inches(4.7), Inches(10), Inches(1.8), size=14, color=BODY_TEXT)

# =============================================
# SLIDE 10: Managing Experiments
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "Managing Experiments", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

# Three feature boxes
features = [
    ("Edit Past Results", ACCENT,
     ["Correct sensory ratings if a panel needs re-evaluation",
      "System warns that later experiments may be invalidated",
      "Utility scores are recalculated automatically"]),
    ("Rewind", RGBColor(0x29, 0x80, 0xB9),
     ["Roll back to any past experiment",
      "Discard subsequent iterations that are no longer valid",
      "Original state is archived automatically (nothing lost)"]),
    ("Delete", RGBColor(0xE6, 0x7E, 0x22),
     ["Remove individual outlier experiments",
      "History re-indexes automatically",
      "Useful for discarding failed trials"]),
]

for i, (title, color, bullets) in enumerate(features):
    x = Inches(0.8 + i * 4.2)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.8), Inches(3.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = color

    add_title_text(slide, title, x + Inches(0.2), Inches(1.7), Inches(3.4), Inches(0.5),
                   size=22, color=color, bold=True)
    add_bullet_list(slide, bullets, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(2.5), size=15)

# =============================================
# SLIDE 11: Rewind Workflow
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "Workflow: Correcting a Past Evaluation", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

add_body_text(slide, "Scenario: At iteration 6, you realize iteration 4's sensory panel was flawed.",
              Inches(1.0), Inches(1.5), Inches(10), Inches(0.5), size=20, color=DARK_BLUE)

steps = [
    ("1", "Edit experiment #4 with\ncorrected sensory ratings"),
    ("2", "Rewind to experiment #4\n(discards #5 and #6)"),
    ("3", "System re-optimizes\nfrom the corrected data"),
    ("4", "New suggestion for #5\nreflects updated ratings"),
]

for i, (num, desc) in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.8)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.8), y, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_body_text(slide, desc, x + Inches(0.1), y + Inches(1.0), Inches(2.2), Inches(1),
                  size=16, color=BODY_TEXT)

    if i < 3:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(2.5), y + Inches(0.15), Inches(0.5), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
        arrow.line.fill.background()

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.2), Inches(11), Inches(1.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
shape.line.color.rgb = ACCENT
txf = shape.text_frame
txf.word_wrap = True
txf.paragraphs[0].text = "The archive preserves all prior work for reference - nothing is permanently lost."
txf.paragraphs[0].font.size = Pt(18)
txf.paragraphs[0].font.bold = True
txf.paragraphs[0].font.color.rgb = DARK_BLUE
txf.paragraphs[0].alignment = PP_ALIGN.CENTER

# =============================================
# SLIDE 12: Technology Stack
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "Technology Stack", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

add_table(slide,
    ["Component", "Technology", "Purpose"],
    [
        ["Optimization Engine", "BoTorch + GPyTorch (Meta AI)", "State-of-the-art Bayesian optimization"],
        ["Surrogate Model", "Gaussian Process (SingleTaskGP)", "Probabilistic model of ingredient-outcome relationships"],
        ["Acquisition Function", "qLogNoisyExpectedImprovement", "Smart selection of next experiments"],
        ["User Interface", "Streamlit", "Interactive web application"],
        ["Language", "Python + PyTorch", "Robust scientific computing foundation"],
    ],
    Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5)
)

# =============================================
# SLIDE 13: Key Benefits
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "Key Benefits", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

benefits = [
    ("Fewer Experiments", "Reach optimal formulation in 10-15 iterations\nvs. 50-100+ with traditional methods"),
    ("Multi-Objective", "Optimize multiple sensory attributes\nsimultaneously with flexible goals"),
    ("Constraint-Aware", "Every suggestion respects cost, nutrition,\nand mass constraints automatically"),
    ("Batch Suggestions", "Run multiple experiments per iteration\nfor efficient parallel lab work"),
    ("Full Audit Trail", "Edit, delete, and rewind capabilities\nwith automatic archiving"),
    ("No ML Expertise", "Intuitive web interface - no coding or\nmachine learning knowledge required"),
]

for i, (title, desc) in enumerate(benefits):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.2)
    y = Inches(1.6 + row * 2.8)

    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.8), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = ACCENT

    add_title_text(slide, title, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.5),
                   size=20, color=ACCENT, bold=True)
    add_body_text(slide, desc, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(1.2),
                  size=15, color=BODY_TEXT)

# =============================================
# SLIDE 14: Next Steps
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_accent_bar(slide)
add_title_text(slide, "Next Steps", Inches(0.8), Inches(0.5), Inches(11), Inches(0.8), size=36)
add_footer(slide)

steps = [
    "1.  Define your ingredient list and allowed ranges",
    "2.  Set sensory objectives (what to maximize, minimize, or target)",
    "3.  Configure constraints (cost, nutrition, total mass)",
    "4.  Run the first batch of 3-5 experiments",
    "5.  Evaluate with your sensory panel and input ratings",
    "6.  Iterate - the system improves with every data point",
]
add_bullet_list(slide, steps, Inches(1.5), Inches(1.8), Inches(8), Inches(4), size=22)

# =============================================
# SLIDE 15: Thank You
# =============================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.12))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

add_title_text(slide, "Thank You",
               Inches(1), Inches(2.2), Inches(11), Inches(1.2),
               size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_title_text(slide, "Questions?",
               Inches(1), Inches(3.5), Inches(11), Inches(0.8),
               size=32, color=RGBColor(0xBD, 0xC3, 0xC7), bold=False, alignment=PP_ALIGN.CENTER)
add_title_text(slide, "Formulation Assistant - Smarter experiments, faster results.",
               Inches(1), Inches(5.0), Inches(11), Inches(0.6),
               size=18, color=RGBColor(0x95, 0xA5, 0xA6), bold=False, alignment=PP_ALIGN.CENTER)

# --- Save ---
output_path = "Formulation_Assistant_Overview.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
